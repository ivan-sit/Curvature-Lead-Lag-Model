#!/usr/bin/env python3
"""Generate the ~20-minute talk deck (PPTX) for the Ricci-curvature lead-lag study.

Clean, concise slides — few points each — following the arc:
  why Ricci curvature  ->  hypothesis  ->  what I did  ->  structural results  ->  implications

Run:  python scripts/make_slides.py
Out:  slides/Ricci_Curvature_LeadLag.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---- palette -------------------------------------------------------------- #
INK = RGBColor(0x1F, 0x29, 0x37)      # near-black slate (titles)
BODY = RGBColor(0x37, 0x41, 0x51)     # dark gray (body text)
ACCENT = RGBColor(0x25, 0x63, 0xEB)   # blue (rules, accents)
MUTE = RGBColor(0x9C, 0xA3, 0xAF)     # light gray (footer)
PALE = RGBColor(0xEF, 0xF2, 0xFB)     # pale blue (table zebra)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x15, 0x9E, 0x5B)    # confirmed
AMBER = RGBColor(0xD9, 0x77, 0x06)    # honest null / not supported
PALE_G = RGBColor(0xEC, 0xFB, 0xF2)   # pale green panel
PALE_A = RGBColor(0xFD, 0xF4, 0xE6)   # pale amber panel
FONT = "Calibri"

# ---- geometry (16:9) ------------------------------------------------------ #
SW, SH = Inches(13.333), Inches(7.5)
ML = Inches(0.85)             # left margin
CW = Inches(11.6)             # content width
TITLE_TOP = Inches(0.7)


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _rule(slide, top, left=ML, width=Inches(2.0), color=ACCENT, h=Pt(3)):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _footer(slide, idx, section):
    _, tf = _box(slide, ML, Inches(6.95), CW, Inches(0.4))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"{section}"
    _set(r, 11, MUTE)
    # page number right-aligned
    _, tf2 = _box(slide, Inches(11.8), Inches(6.95), Inches(1.4), Inches(0.4))
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = str(idx)
    _set(r2, 11, MUTE)


def title_slide(prs, title, subtitle, meta):
    s = _blank(prs)
    _rule(s, Inches(2.55), left=ML, width=Inches(2.4), color=ACCENT, h=Pt(4))
    _, tf = _box(s, ML, Inches(2.75), CW, Inches(2.2))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set(r, 40, INK, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(14)
    r2 = p2.add_run(); r2.text = subtitle
    _set(r2, 22, ACCENT, italic=True)
    _, tf3 = _box(s, ML, Inches(5.35), CW, Inches(1.2))
    for i, line in enumerate(meta):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        r = p.add_run(); r.text = line
        _set(r, 15, BODY)
    return s


def content_slide(prs, idx, section, title, bullets, lead=None, closer=None,
                  closer_label="Research question"):
    """bullets: list of str (level 0) or (str, level) tuples.

    ``closer`` (optional): a highlighted call-out box near the bottom — used to
    land a research question, a horizon summary, or a one-line punch. Its bold
    blue tag is ``closer_label``.
    """
    s = _blank(prs)
    _, tf = _box(s, ML, TITLE_TOP, CW, Inches(0.9))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set(r, 30, INK, bold=True)
    _rule(s, Inches(1.55))
    top = Inches(1.95)
    if lead:
        _, lf = _box(s, ML, top, CW, Inches(0.8))
        pp = lf.paragraphs[0]
        rr = pp.add_run(); rr.text = lead
        _set(rr, 19, ACCENT, italic=True)
        top = Inches(2.75)
    _, bf = _box(s, ML, top, CW, Inches(3.0 if closer else 4.0))
    for i, b in enumerate(bullets):
        text, level = (b if isinstance(b, tuple) else (b, 0))
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.level = level
        p.space_after = Pt(12 if level == 0 else 6)
        bullet = "•  " if level == 0 else "–  "
        r = p.add_run(); r.text = bullet + text
        _set(r, 21 if level == 0 else 18, BODY if level == 0 else MUTE,
             bold=False)
    if closer:
        from pptx.enum.shapes import MSO_SHAPE
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ML, Inches(5.45),
                                 CW, Inches(1.15))
        box.fill.solid(); box.fill.fore_color.rgb = PALE
        box.line.color.rgb = ACCENT; box.line.width = Pt(1.25)
        box.shadow.inherit = False
        ctf = box.text_frame; ctf.word_wrap = True
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = Inches(0.25); ctf.margin_right = Inches(0.25)
        cp = ctf.paragraphs[0]
        lbl = cp.add_run(); lbl.text = closer_label + "   "
        _set(lbl, 14, ACCENT, bold=True)
        rr = cp.add_run(); rr.text = closer
        _set(rr, 17, INK, bold=True)
    _footer(s, idx, section)
    return s


def _fill_cell(cell, val, size, color, bold, align):
    """Write possibly multi-line text into a table cell (split on '\\n')."""
    tf = cell.text_frame
    for k, line in enumerate(str(val).split("\n")):
        para = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run(); run.text = line
        _set(run, size, color, bold=bold)


def table_slide(prs, idx, section, title, headers, rows, note=None, lead=None,
                col_widths=None, highlight_rows=(), align=None):
    s = _blank(prs)
    _, tf = _box(s, ML, TITLE_TOP, CW, Inches(0.9))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set(r, 30, INK, bold=True)
    _rule(s, Inches(1.55))
    top = Inches(1.95)
    if lead:
        _, lf = _box(s, ML, top, CW, Inches(0.7))
        pp = lf.paragraphs[0]
        rr = pp.add_run(); rr.text = lead
        _set(rr, 18, ACCENT, italic=True)
        top = Inches(2.6)
    nr, nc = len(rows) + 1, len(headers)
    gtbl = s.shapes.add_table(nr, nc, ML, top, CW, Inches(0.5 * nr)).table
    if col_widths:
        for j, w in enumerate(col_widths):
            gtbl.columns[j].width = w
    amap = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
    aligns = [amap[a] for a in align] if align else \
        [PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * (nc - 1)
    # header
    for j, h in enumerate(headers):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = ACCENT
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        _fill_cell(c, h, 15, WHITE, True, aligns[j])
    # body
    for i, row in enumerate(rows, start=1):
        hl = (i - 1) in highlight_rows
        for j, val in enumerate(row):
            c = gtbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = PALE if (hl or i % 2 == 0) else WHITE
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            _fill_cell(c, val, 14, INK if hl else BODY, hl, aligns[j])
    if note:
        _, nf = _box(s, ML, Inches(6.25), CW, Inches(0.7))
        pp = nf.paragraphs[0]
        rr = pp.add_run(); rr.text = note
        _set(rr, 14, MUTE, italic=True)
    _footer(s, idx, section)
    return s


def _flowbox(slide, x, y, w, h, title, sub=None, accent=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = ACCENT if accent else PALE
    shp.line.color.rgb = ACCENT; shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    _set(r, 14, WHITE if accent else INK, bold=True)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(2)
        r2 = p2.add_run(); r2.text = sub
        _set(r2, 10, RGBColor(0xE5, 0xED, 0xFF) if accent else MUTE)


def _arrow(slide, kind, x, y, w, h):
    shp = slide.shapes.add_shape(kind, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = ACCENT
    shp.line.fill.background(); shp.shadow.inherit = False


def pipeline_slide(prs, idx, section, title, stages, lead=None):
    """Six-stage pipeline as a snake flow chart (3 top L->R, 3 bottom R->L)."""
    s = _blank(prs)
    _, tf = _box(s, ML, TITLE_TOP, CW, Inches(0.9))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set(r, 30, INK, bold=True)
    _rule(s, Inches(1.55))
    if lead:
        _, lf = _box(s, ML, Inches(1.75), CW, Inches(0.6))
        rr = lf.paragraphs[0].add_run(); rr.text = lead
        _set(rr, 18, ACCENT, italic=True)
    bw, bh = Inches(3.3), Inches(1.1)
    xs = [ML, Inches(4.6), Inches(8.35)]
    ytop, ybot = Inches(2.55), Inches(4.75)
    ah = Inches(0.34)
    # top row, left -> right
    for k in range(3):
        _flowbox(s, xs[k], ytop, bw, bh, stages[k][0], stages[k][1])
    _arrow(s, MSO_SHAPE.RIGHT_ARROW, Inches(4.18), ytop + Inches(0.38), Inches(0.42), ah)
    _arrow(s, MSO_SHAPE.RIGHT_ARROW, Inches(7.93), ytop + Inches(0.38), Inches(0.42), ah)
    # down on the right
    _arrow(s, MSO_SHAPE.DOWN_ARROW, Inches(9.83), Inches(3.7), ah, Inches(1.0))
    # bottom row, right -> left (so the flow continues 4 -> 5 -> 6)
    bottom_x = [Inches(8.35), Inches(4.6), ML]
    for k in range(3):
        _flowbox(s, bottom_x[k], ybot, bw, bh, stages[3 + k][0], stages[3 + k][1],
                 accent=(3 + k == 5))
    _arrow(s, MSO_SHAPE.LEFT_ARROW, Inches(7.93), ybot + Inches(0.38), Inches(0.42), ah)
    _arrow(s, MSO_SHAPE.LEFT_ARROW, Inches(4.18), ybot + Inches(0.38), Inches(0.42), ah)
    _footer(s, idx, section)
    return s


def verdict_slide(prs, idx, section, title, left, right, lead=None):
    """Two-column hypothesis verdict. left/right = dict(label, verdict, color,
    pale, points[list])."""
    s = _blank(prs)
    _, tf = _box(s, ML, TITLE_TOP, CW, Inches(0.9))
    r = tf.paragraphs[0].add_run(); r.text = title
    _set(r, 30, INK, bold=True)
    _rule(s, Inches(1.55))
    if lead:
        _, lf = _box(s, ML, Inches(1.75), CW, Inches(0.5))
        rr = lf.paragraphs[0].add_run(); rr.text = lead
        _set(rr, 17, ACCENT, italic=True)
    colw = Inches(5.55)
    xs = [ML, Inches(6.9)]
    top = Inches(2.45)
    for col, x in zip((left, right), xs):
        # header band
        band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, colw, Inches(0.7))
        band.fill.solid(); band.fill.fore_color.rgb = col["color"]
        band.line.fill.background(); band.shadow.inherit = False
        btf = band.text_frame; btf.word_wrap = True
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = col["label"]
        _set(br, 16, WHITE, bold=True)
        # body panel
        panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.3),
                                   colw, Inches(3.2))
        panel.fill.solid(); panel.fill.fore_color.rgb = col["pale"]
        panel.line.color.rgb = col["color"]; panel.line.width = Pt(1)
        panel.shadow.inherit = False
        ptf = panel.text_frame; ptf.word_wrap = True
        ptf.margin_left = Inches(0.3); ptf.margin_right = Inches(0.3)
        ptf.margin_top = Inches(0.25)
        vp = ptf.paragraphs[0]
        vr = vp.add_run(); vr.text = col["verdict"]
        _set(vr, 22, col["color"], bold=True)
        for pt in col["points"]:
            pp = ptf.add_paragraph(); pp.space_before = Pt(10)
            rr = pp.add_run(); rr.text = "•  " + pt
            _set(rr, 15, BODY)
    _footer(s, idx, section)
    return s


def section_slide(prs, idx, kicker, big):
    s = _blank(prs)
    _rule(s, Inches(3.0), left=ML, width=Inches(1.8), color=ACCENT, h=Pt(4))
    _, tf = _box(s, ML, Inches(3.2), CW, Inches(1.6))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = kicker
    _set(r, 16, ACCENT, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(6)
    r2 = p2.add_run(); r2.text = big
    _set(r2, 34, INK, bold=True)
    return s


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    n = 0

    def nx():
        nonlocal n; n += 1; return n

    # 1 — title (+ acronym CRISP)
    title_slide(
        prs,
        "CRISP: Curvature of Residualized, Signed Lead-Lag Pairs",
        "Directed Forman-Ricci curvature as a structural lens on US equity information flow",
        ["Ivan Sit", "UCLA Math 285J — Agentic AI for Autonomous Research in Quant Finance",
         "Advisor: Prof. Mihai Cucuringu   ·   Mid-project seminar, June 2026   ·   Target: ICAIF 2026"],
    )

    # 2 — what is curvature
    content_slide(
        prs, nx(), "Motivation", "Curvature, on a graph",
        ["Ricci curvature measures how a space bends — and how diffusion flows through it.",
         "On a graph it lives on each edge and has one intuitive meaning:",
         ("Negative = a bridge: connects separated regions, structurally isolated.", 1),
         ("Positive = an interior link: sits inside a dense cluster, many paths route around it.", 1),
         "A parameter-free, local number with a global meaning."],
        lead="Is this relationship a fragile bridge, or a redundant interior link?",
    )

    # 3 — why I got into it
    content_slide(
        prs, nx(), "Motivation", "Why I got into Ricci curvature",
        ["Sandhu, Georgiou & Tannenbaum (Science Advances, 2016):",
         ("Ollivier-Ricci curvature on equity correlation networks rises before crashes.", 1),
         ("Curvature as a systemic-fragility signal — geometry, not statistics.", 1),
         "It reframes a market as a shape whose bottlenecks you can measure.",
         "That raised the question this project asks."],
    )

    # 4 — the one opening -> research question
    content_slide(
        prs, nx(), "Motivation", "The opening they left",
        ["Sandhu et al. built the network on an undirected correlation graph.",
         "But markets have direction — some names lead, others lag.",
         ("The symmetric correlation graph erases that asymmetry by construction.", 1)],
        lead="One thing their construction could not see: direction.",
        closer="Does Ricci curvature on the directed, residualized lead-lag network reveal "
               "structure the undirected correlation view cannot?",
    )

    # 5 — two hypotheses
    table_slide(
        prs, nx(), "Hypothesis", "Two hypotheses",
        ["H1 — Structural  (primary)", "H2 — Predictive  (secondary)"],
        [["Negatively-curved pairs in the directed, residualized lead-lag network are "
          "structurally distinct — bridges that correlation, degree, and undirected "
          "curvature do not recover.",
          "Those same pairs also carry out-of-sample directional information: today's "
          "leader predicts tomorrow's lagger, beyond the baselines."]],
        lead="Both falsifiable. The paper's claim rests on H1; H2 is exploratory.",
        note="We answer both at the end.",
        col_widths=[Inches(5.7), Inches(5.7)],
        align=["l", "l"],
    )

    # 5b — related work & contribution
    table_slide(
        prs, nx(), "Related work", "Related literature & contribution",
        ["Prior work", "What it did", "What CRISP adds"],
        [["Sandhu et al. 2016", "Ollivier-Ricci on undirected correlation nets",
          "directed lead-lag + Forman"],
         ["Forman 2003; Samal 2018", "Forman / augmented curvature on graphs",
          "directed, weighted, residualized"],
         ["Bennett–Cucuringu–Reinert 2022", "signed lead-lag estimator + significance",
          "adopted as the edge weight"],
         ["Tian–Lubberts–Weber 2025", "line-graph curvature clustering (undirected)",
          "pair-space view on L(G)"]],
        lead="Gap: nobody has put directed curvature on a residualized lead-lag graph.",
        note="Contribution: the directed-vs-undirected and Forman-vs-Ollivier contrasts, "
             "on factor-residualized returns — a structural extension of Sandhu.",
        col_widths=[Inches(3.6), Inches(4.6), Inches(3.4)],
        align=["l", "l", "l"],
    )

    # 6 — what I did, overview (flow chart)
    pipeline_slide(
        prs, nx(), "What I did", "The pipeline",
        [("Residualize returns", "remove market + sector"),
         ("Directed lead-lag network", "BCR signed statistic"),
         ("Four curvature objects", "plain → weighted aug. directed"),
         ("Line graph  L(G)", "curvature pair-communities"),
         ("Validation cascade", "curvature ≠ corr / degree?"),
         ("Structurally isolated pairs", "the structural output")],
        lead="Six stages, fully automated.",
    )

    # 7 — residualized lead-lag (+ explicit horizons)
    content_slide(
        prs, nx(), "What I did", "Step 1 — Residualized, directed network",
        ["Residualize first: on raw returns, market + sector dominate the graph —",
         ("otherwise curvature just rediscovers GICS sectors (instructor requirement).", 1),
         "Edge weight = BCR signed lead-lag:  w(i→j) = ρ_ij(τ*) − ρ_ji(τ*).",
         ("τ* = the lag of peak signal, picked per pair; the candidate horizons are on "
          "the next slide.", 1),
         "Two time horizons: daily close-to-close, and intraday 30-min.",
         ("Intraday uses a within-day estimator — lag pairs never cross the overnight gap.", 1)],
    )

    # 8 — four objects (with plain-language "what it measures")
    table_slide(
        prs, nx(), "What I did", "Step 2 — Four curvature objects",
        ["Object", "What it measures", "Role"],
        [["Plain directed Forman\n4 − deg_in − deg_out",
          "only the two endpoints' link counts — pure connectivity",
          "Degree BASELINE\n(no higher-order signal)"],
         ["Triangle-augmented\nF# = F + 3m",
          "+3 for every triangle the edge sits inside",
          "Higher-order:\nisolates the 3m term"],
         ["Weighted augmented\ndirected Forman",
          "lead-lag strength + triangles + edge direction",
          "MAIN object"],
         ["Ollivier-Ricci\n(optimal transport)",
          "cost to morph one node's neighborhood into the other's",
          "Robustness / contrast\n(Sandhu's)"]],
        lead="One ablation, from pure degree to full directed curvature.",
        note="Plain Forman is an exact function of endpoint degrees — the baseline the rest must beat.",
        col_widths=[Inches(3.6), Inches(4.7), Inches(3.3)],
        align=["l", "l", "l"],
        highlight_rows=(2,),
    )

    # 9 — cascade
    content_slide(
        prs, nx(), "What I did", "Step 3 — The validation cascade",
        ["Goal: prove curvature is not a re-skin of correlation, degree, or sector.",
         "Spearman(F, |ρ|) — must be small.",
         "Top-K Jaccard vs correlation-ranked pairs — must be small.",
         "Degree-preserving configuration-model null (z-scores).",
         "Residual orthogonalization — regress out degree & |ρ|, keep the rest.",
         "Benjamini–Hochberg + strict train / validate / test split."],
        lead="This cascade is the methodological core.",
    )

    # 10 — data (what each source is FOR)
    table_slide(
        prs, nx(), "What I did", "Data — all institutional (WRDS)",
        ["Source", "What it is", "What it's for"],
        [["CRSP", "daily S&P 500 returns, 2000–2024\n(survivorship-correct)",
          "structural cascade + 25-yr predictive walk-forward"],
         ["Compustat", "GICS sector classifications",
          "residualize returns; sector-boundary tests"],
         ["TAQ", "30-min intraday bars, 2019 only\n(~155 large-caps)",
          "high-frequency cross-check (lead-lag is strongest)"]],
        lead="Every network is built on factor-residualized returns.",
        note="Why TAQ is 2019 only: 25 years of 30-min TAQ ≈ 100+ hours of pulls (TAQ is ~15B "
             "rows/year). Daily covers the full 2000–2024; intraday 2019 is the cross-check.",
        col_widths=[Inches(2.0), Inches(5.4), Inches(4.2)],
        align=["l", "l", "l"],
        highlight_rows=(2,),
    )

    # 11 — the two regimes / time horizons we ran
    table_slide(
        prs, nx(), "What I did", "Two regimes — two time horizons",
        ["", "Intraday  (headline)", "Daily  (robustness)"],
        [["Frequency", "30-min bars", "close-to-close daily"],
         ["Lead-lag horizon  τ", "1–3 bars  =  30 / 60 / 90 min", "1–5 trading days"],
         ["Universe", "~155 large-caps", "S&P 500 (point-in-time)"],
         ["Window", "full-year 2019", "2000–2024"]],
        lead="The same method, run at two horizons — no weekly or monthly.",
        note="Coverage: daily spans the full 2000–2024; intraday is 2019 only (25 yrs of 30-min "
             "TAQ ≈ 100+ h of pulls), so it serves as the high-frequency cross-check.",
        col_widths=[Inches(2.9), Inches(4.5), Inches(4.2)],
        align=["l", "c", "c"],
        highlight_rows=(1,),
    )

    # 12 — section: results
    section_slide(prs, nx(), "RESULTS", "Two hypotheses, two answers.")

    # 13 — H1 result 1: distinctness (BOTH horizons)
    table_slide(
        prs, nx(), "Results — H1", "Curvature is not correlation",
        ["Test", "Intraday\n2019", "Daily\n2000–24", "What it means"],
        [["Top-K Jaccard vs correlation", "≈ 0.0", "0.0", "near-disjoint pair sets"],
         ["Spearman(F, |ρ|)", "0.18", "0.07", "barely tracks |ρ|  (≪ 0.8)"]],
        lead="The pairs curvature flags are not the pairs correlation flags — at both horizons.",
        note="Significance: curvature is not a re-skin of correlation. Holds on the 2019 intraday "
             "graph AND the full 2000–2024 daily span. This distinctness IS the headline result.",
        col_widths=[Inches(3.7), Inches(1.9), Inches(2.0), Inches(4.0)],
        align=["l", "c", "c", "l"],
        highlight_rows=(0, 1),
    )

    # 13b — the four objects compared (the ablation grid)
    table_slide(
        prs, nx(), "Results — H1", "The four objects, compared",
        ["Object", "Spearman", "R² on degree", "mean IC", "compute"],
        [["Plain directed Forman  (baseline)", "0.09", "1.00", "+0.014", "0.03 s"],
         ["Weighted Forman", "0.08", "0.97", "+0.007", "0.03 s"],
         ["Weighted augmented directed  (MAIN)", "0.18", "0.56", "+0.005", "0.03 s"],
         ["Ollivier-Ricci", "0.24", "0.14", "−0.006", "15.9 s"]],
        lead="All four pick pairs disjoint from correlation (Jaccard ≈ 0).",
        note="Best tradeoff = weighted-augmented Forman: real higher-order signal (44% non-degree) "
             "at ~0 cost. Ollivier carries the most non-degree signal but at ~600× the compute and "
             "no IC gain. None predicts (IC ≈ 0).",
        col_widths=[Inches(4.4), Inches(1.7), Inches(2.1), Inches(1.7), Inches(1.7)],
        align=["l", "c", "c", "c", "c"],
        highlight_rows=(2,),
    )

    # 14 — H1 result 2: degree ablation (BOTH horizons)
    table_slide(
        prs, nx(), "Results — H1", "A clean degree ablation",
        ["Object  (R² on degree)", "Intraday\n2019", "Daily\n2000–24", "What it means"],
        [["Plain Forman", "1.00", "1.00", "100% degree — pure baseline, by design"],
         ["Augmented Forman", "0.56", "0.18", "44–82% is NOT degree → higher-order signal"]],
        lead="Does curvature carry anything beyond connectivity? — checked at both horizons.",
        note="Significance: plain Forman = pure degree (calibration ✓); the augmented object adds "
             "real higher-order structure at both horizons — even more on daily (~82% non-degree).",
        col_widths=[Inches(3.7), Inches(1.9), Inches(2.0), Inches(4.0)],
        align=["l", "c", "c", "l"],
        highlight_rows=(1,),
    )

    # 14 — H1 result 3: triangle sparsity
    content_slide(
        prs, nx(), "Results — H1", "The structure is cross-sector",
        ["Almost all of the network's triangle structure spans sectors —",
         ("only ~1% of triangles sit within a single GICS sector.", 1),
         "So the relationships curvature highlights are cross-sector links — exactly what "
         "a sector- or correlation-based grouping would not surface.",
         "Honest caveat: with so few within-group triangles, the community step can't "
         "cleanly carve out clusters."],
        lead="Where do the pairs curvature picks actually live?",
    )

    # 15 — H2 result: the predictive test (IC), BOTH horizons
    table_slide(
        prs, nx(), "Results — H2", "The predictive test — both horizons",
        ["Method  (mean OOS IC)", "Intraday 30-min\n2019", "Daily close-to-close\n2000–2024"],
        [["Curvature (aug, directed)", "+0.005", "−0.011  (CI < 0)"],
         ["Undirected Forman", "+0.013", "+0.002"],
         ["Correlation", "−0.003", "+0.023  (CI > 0)"],
         ["Random", "−0.015", "+0.007"]],
        lead="Tested at both horizons — intraday (2019) and daily (25-yr walk-forward).",
        note="Intraday: Forman leads but every CI spans zero. Daily (21 windows): curvature "
             "is significantly NEGATIVE and correlation does best. No predictive edge for curvature.",
        col_widths=[Inches(4.6), Inches(3.3), Inches(3.7)],
        align=["l", "c", "c"],
        highlight_rows=(0,),
    )

    # 16 — verdict: respond to BOTH hypotheses
    verdict_slide(
        prs, nx(), "Verdict", "Did the two hypotheses hold?",
        left={
            "label": "H1 — Structural distinctness", "color": GREEN, "pale": PALE_G,
            "verdict": "✓  CONFIRMED",
            "points": ["Jaccard 0, Spearman 0.07–0.18 vs correlation",
                       "augmented Forman: 44–82% non-degree signal",
                       "robust across 2019 intraday & 25-yr daily"],
        },
        right={
            "label": "H2 — Predictive edge", "color": AMBER, "pale": PALE_A,
            "verdict": "✗  NOT SUPPORTED",
            "points": ["Intraday 2019: Forman leads, but CIs span zero",
                       "Daily 25-yr: curvature negative, correlation best",
                       "no predictive edge; the claim rests on H1"],
        },
        lead="The paper stands on H1. H2 is reported transparently — across 25 years.",
    )

    # 18 — robustness & ablations
    content_slide(
        prs, nx(), "Robustness", "Robustness & ablations",
        ["Four-object ablation: plain (degree baseline) → augmented → weighted-directed → Ollivier.",
         "Two horizons: intraday 30-min (2019) and daily close-to-close (2000–2024).",
         "25-year walk-forward: 21 rolling windows, not a single train/test split.",
         "Five baselines: correlation, cointegration, random, undirected Forman, Ollivier.",
         "Triangle convention (common-neighbor vs strict-cyclic) settled by the data."],
        lead="The result is stress-tested, not cherry-picked.",
    )

    # 19 — agentic AI workflow (single, concise)
    content_slide(
        prs, nx(), "Agentic AI", "Agentic AI workflow",
        ["Helped me research the idea — surveying the curvature and lead-lag literature, "
         "and coding the pipeline.",
         "Its core role: linking curvature (graph geometry) to directed lead-lag networks "
         "(the finance side).",
         "Framing and judgment stayed human."],
        lead="What agentic AI did here.",
    )

    # 20 — limitations
    content_slide(
        prs, nx(), "Limitations", "Limitations",
        ["Intraday is 2019 only (TAQ cost) — not a 25-year intraday panel.",
         "Triangle-sparse network limits AFRC community separability (Fesser–Weber–Lambiotte).",
         "Directed line-graph & directed curvature-gap are theoretically open — undirected reduction used.",
         "Predictive IC evaluated at fixed k = 20; not pre-registered.",
         "Daily universe is survivor-biased within each window (continuous-presence filter)."],
        lead="Stated plainly — these shape what we can and can't claim.",
    )

    # 22 — what's next (active work + submission)
    content_slide(
        prs, nx(), "What's next", "What's next",
        ["Multi-regime intraday robustness — pulling 2008 / 2015 / 2020 to join 2019 "
         "(crash, calm, COVID).",
         "Agentic propose→test→reject orchestrator (built) — logs every accept/reject; "
         "the residualization ladder (market / sector / PCA) runs through it.",
         "Directed line-graph & directed curvature-gap theory (open; possible Weber collaboration).",
         "Extend to an 8-page ACM sigconf — ICAIF 2026, deadline Aug 2."],
        lead="Underway now.",
    )

    # 23 — project status overview
    table_slide(
        prs, nx(), "Status", "Project Status Overview",
        ["✓  Completed", "⟳  In progress", "▢  Remaining before final"],
        [["Full pipeline + 52 tests\n"
          "Four curvature objects\n"
          "Structural cascade (H1) ✓\n"
          "Predictive test (H2): 2 horizons, 25 yrs\n"
          "Agentic orchestrator (propose→test→reject)\n"
          "Residualization-ladder ablation\n"
          "Report + deck + talk script",
          "Multi-regime intraday pull\n"
          "(2008 / 2015 / 2020)\n"
          "LaTeX/Overleaf main.tex\n"
          "Bibliography → 20–25 refs",
          "Directed line-graph theory\n"
          "Final 8-page ICAIF write-up\n"
          "Paper-quality figures"]],
        lead="Mid-project status — ~10 days to the final write-up.",
        col_widths=[Inches(4.1), Inches(3.5), Inches(4.0)],
        align=["l", "l", "l"],
    )

    out = Path("slides"); out.mkdir(exist_ok=True)
    path = out / "Ricci_Curvature_LeadLag.pptx"
    prs.save(str(path))
    print(f"saved {path}  ({n} slides)")


if __name__ == "__main__":
    build()
